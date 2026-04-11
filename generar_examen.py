#!/usr/bin/env python3
"""
generar_examen.py — Genera un examen Cambridge C1/B2 en PPTX a partir de JSONs.

Uso:
  python generar_examen.py <ruta_test> [plantilla.pptx]

Ejemplos:
  python generar_examen.py Nivel/C1/Exams/Test1
  python generar_examen.py Nivel/C1/Exams/Test1 mi_plantilla.pptx
"""

import copy
import json
import os
import re
import sys

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor

# Namespace DrawingML
_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _tag(local):
    return f"{{{_NS_A}}}{local}"


def _save_para_template(tf):
    """
    Devuelve una copia profunda del primer párrafo del text frame
    para usarla como plantilla de formato al reconstruir el contenido.
    """
    paras = tf._txBody.findall(_tag("p"))
    if paras:
        return copy.deepcopy(paras[0])
    return None


def _build_run(rPr_template, text):
    """Crea un elemento <a:r> con el rPr copiado de la plantilla y el texto dado."""
    r = etree.Element(_tag("r"))
    if rPr_template is not None:
        r.append(copy.deepcopy(rPr_template))
    t = etree.SubElement(r, _tag("t"))
    t.text = text
    # xml:space="preserve" para mantener espacios iniciales/finales
    t.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")
    return r


def _build_para_from_template(template_p, text):
    """
    Construye un <a:p> nuevo con el texto dado, conservando pPr y rPr
    del párrafo plantilla.
    """
    new_p = etree.Element(_tag("p"))

    # Copiar pPr si existe
    pPr = template_p.find(_tag("pPr"))
    if pPr is not None:
        new_p.append(copy.deepcopy(pPr))

    # Obtener rPr del primer run de la plantilla
    first_r = template_p.find(_tag("r"))
    rPr = None
    if first_r is not None:
        rPr_elem = first_r.find(_tag("rPr"))
        if rPr_elem is not None:
            rPr = rPr_elem

    new_p.append(_build_run(rPr, text))
    return new_p


def _replace_txBody_text(tf, lines, template_p):
    """
    Reemplaza el contenido del text frame con las líneas dadas, conservando
    el formato del párrafo/run plantilla.
    """
    txBody = tf._txBody
    # Eliminar todos los párrafos actuales
    for p in txBody.findall(_tag("p")):
        txBody.remove(p)
    # Añadir un párrafo por línea
    for line in lines:
        if template_p is not None:
            new_p = _build_para_from_template(template_p, line)
        else:
            new_p = etree.Element(_tag("p"))
            r = etree.SubElement(new_p, _tag("r"))
            t = etree.SubElement(r, _tag("t"))
            t.text = line
        txBody.append(new_p)

TEMPLATE_DEFAULT = "examen_plantilla.pptx"

# Indicador de hueco en el texto (puntos de relleno)
GAP_DOTS = "…………"


# ─── helpers silenciosos (sin avisos ni warnings) ────────────────────────────

def get_shape(slide, name):
    """Devuelve la forma con ese nombre, o None si no existe."""
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def get_table(slide, name):
    """Devuelve la tabla (objeto Table) con ese nombre, o None si no existe."""
    shape = get_shape(slide, name)
    if shape is None or not shape.has_table:
        return None
    return shape.table


def set_text(slide, name, text):
    """
    Escribe texto plano en una forma conservando el estilo original del cuadro
    de texto (fuente, tamaño, color, alineación, etc.).
    No hace nada si la forma no existe.
    """
    shape = get_shape(slide, name)
    if shape is None or not shape.has_text_frame:
        return
    tf = shape.text_frame
    template_p = _save_para_template(tf)
    lines = text.split("\n")
    _replace_txBody_text(tf, lines, template_p)


def set_shape_text_rich(slide, name, runs):
    """
    Escribe texto enriquecido en una forma conservando el estilo base del
    cuadro de texto (fuente, tamaño, color…). Solo sobreescribe las
    propiedades explícitamente indicadas en cada run (bold, size).

    runs es una lista de dict:
        {'text': str, 'bold': bool, 'size': int (opcional, en puntos)}

    No hace nada si la forma no existe o no tiene marco de texto.
    """
    shape = get_shape(slide, name)
    if shape is None or not shape.has_text_frame:
        return
    tf = shape.text_frame
    txBody = tf._txBody

    # Guardar plantilla de formato antes de limpiar
    template_p = _save_para_template(tf)
    first_r_tmpl = template_p.find(_tag("r")) if template_p is not None else None
    rPr_base = None
    if first_r_tmpl is not None:
        rPr_elem = first_r_tmpl.find(_tag("rPr"))
        if rPr_elem is not None:
            rPr_base = rPr_elem

    # Eliminar párrafos existentes
    for p in txBody.findall(_tag("p")):
        txBody.remove(p)

    # Construir el único párrafo con todos los runs
    new_p = etree.Element(_tag("p"))
    if template_p is not None:
        pPr = template_p.find(_tag("pPr"))
        if pPr is not None:
            new_p.append(copy.deepcopy(pPr))

    for run_data in runs:
        raw_text = run_data.get("text", "")
        # Dividir por saltos de línea dentro del run
        parts = raw_text.split("\n")
        for part_idx, part in enumerate(parts):
            r = etree.SubElement(new_p, _tag("r"))
            # rPr: partir de la base y aplicar overrides
            if rPr_base is not None:
                rPr = copy.deepcopy(rPr_base)
            else:
                rPr = etree.Element(_tag("rPr"))
                rPr.set("lang", "en-US")
            if run_data.get("bold", False):
                rPr.set("b", "1")
            else:
                rPr.attrib.pop("b", None)
            if "size" in run_data:
                rPr.set("sz", str(int(run_data["size"] * 100)))
            r.insert(0, rPr)
            t = etree.SubElement(r, _tag("t"))
            t.text = part
            t.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")

            # Si no es el último fragmento, añadir salto de línea (nueva línea)
            if part_idx < len(parts) - 1:
                br = etree.SubElement(new_p, _tag("br"))
                if rPr_base is not None:
                    br.append(copy.deepcopy(rPr_base))

    txBody.append(new_p)


def fill_black_shape(slide, name):
    """Rellena una forma de negro. No hace nada si la forma no existe."""
    shape = get_shape(slide, name)
    if shape is None:
        return
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 0, 0)


def clear_shape_fill(slide, name):
    """Elimina el relleno de una forma (fondo transparente/blanco). Silencioso."""
    shape = get_shape(slide, name)
    if shape is None:
        return
    shape.fill.background()


_FILL_TAGS = {
    _tag("solidFill"), _tag("gradFill"), _tag("pattFill"),
    _tag("blipFill"), _tag("grpFill"), _tag("noFill"),
}


def _force_black_text(tf):
    """Fuerza color negro en todos los runs del text frame."""
    for r in tf._txBody.iter(_tag("r")):
        rPr = r.find(_tag("rPr"))
        if rPr is None:
            rPr = etree.Element(_tag("rPr"))
            r.insert(0, rPr)
        # Eliminar fills de color existentes
        for child in list(rPr):
            if child.tag in _FILL_TAGS:
                rPr.remove(child)
        # Insertar solidFill negro al principio de rPr (posición correcta en DML)
        sf = etree.Element(_tag("solidFill"))
        srgb = etree.SubElement(sf, _tag("srgbClr"))
        srgb.set("val", "000000")
        rPr.insert(0, sf)


def fill_table_cell(table, row, col, text, bold=False, force_black=True):
    """
    Escribe texto en una celda de tabla conservando el estilo original de la celda.
    No hace nada si table es None o si la celda no existe.
    force_black=True fuerza color negro en el texto.
    """
    if table is None:
        return
    try:
        cell = table.cell(row, col)
        tf = cell.text_frame
        template_p = _save_para_template(tf)
        lines = text.split("\n") if text else [""]
        _replace_txBody_text(tf, lines, template_p)
        for para in tf.paragraphs:
            for r in para._p.findall(_tag("r")):
                rPr = r.find(_tag("rPr"))
                if rPr is None:
                    rPr = etree.Element(_tag("rPr"))
                    r.insert(0, rPr)
                if bold:
                    rPr.set("b", "1")
                else:
                    rPr.set("b", "0")
        if force_black:
            _force_black_text(tf)
    except Exception:
        pass


def redistribute_shapes_vertically(slide, shape_names):
    """
    Redistribuye verticalmente las formas indicadas para que no se superpongan,
    dividiendo el área total que ocupan entre ellas de forma equitativa.
    """
    shapes = []
    for name in shape_names:
        s = get_shape(slide, name)
        if s is not None:
            shapes.append(s)
    if len(shapes) < 2:
        return
    shapes.sort(key=lambda s: s.top)
    top_start = shapes[0].top
    bottom_end = shapes[-1].top + shapes[-1].height
    total_height = bottom_end - top_start
    n = len(shapes)
    current_top = top_start
    for i, shape in enumerate(shapes):
        # Distribute evenly, assigning the remaining space to the last shape
        if i < n - 1:
            each_height = round(total_height * (i + 1) / n) - round(total_height * i / n)
        else:
            each_height = bottom_end - current_top
        shape.top = current_top
        shape.height = each_height
        current_top += each_height


# ─── procesado de texto ─────────────────────────────────────────────────────

def format_parrafos(raw_text):
    """Convierte el separador || en salto de línea."""
    return raw_text.replace("||", "\n")


def _strip_embedded_answers(raw_text, correct_map):
    """
    Elimina del texto las palabras correctas que aparecen después de cada
    marcador de hueco (N). Ej.: "(0) which" → "(0)".
    """
    for n, correct in correct_map.items():
        if correct:
            raw_text = re.sub(
                rf"\({re.escape(n)}\)\s+{re.escape(correct)}",
                f"({n})",
                raw_text,
                flags=re.IGNORECASE,
            )
    return raw_text


def text_con_huecos(raw_text, questions=None, example=None):
    """
    Sustituye cada (N) en el texto por (N) GAP_DOTS.
    Si se pasan questions/example, elimina previamente la palabra correcta
    que pueda estar incrustada en el texto justo después del marcador.
    """
    correct_map = {}
    if example:
        n = str(example.get("number", 0))
        correct_map[n] = str(example.get("correct", "")).lower()
    if questions:
        for q in questions:
            n = str(q.get("number", ""))
            c = str(q.get("correct", "")).lower()
            if c:
                correct_map[n] = c
    if correct_map:
        raw_text = _strip_embedded_answers(raw_text, correct_map)
    return re.sub(r"\((\d+)\)", rf"(\1) {GAP_DOTS}", raw_text)


def text_con_huecos_wf(raw_text, questions, example=None):
    """
    Word-formation: sustituye (N) CORRECT_WORD por (N) GAP_DOTS (HINT_WORD)
    donde HINT_WORD es la palabra base en mayúsculas de esa pregunta.
    Elimina la palabra correcta que aparezca tras el marcador en el texto.
    """
    correct_map = {}
    hint_map = {}
    if example:
        n = str(example.get("number", 0))
        correct_map[n] = str(example.get("correct", "")).lower()
        hint_map[n] = str(example.get("word", ""))
    for q in questions:
        n = str(q["number"])
        correct_map[n] = str(q.get("correct", "")).lower()
        hint_map[n] = str(q.get("word", ""))

    # Eliminar palabras correctas incrustadas
    raw_text = _strip_embedded_answers(raw_text, correct_map)

    def repl(match):
        n = match.group(1)
        hint = hint_map.get(n, "")
        if hint:
            return f"({n}) {GAP_DOTS} ({hint})"
        return f"({n}) {GAP_DOTS}"

    return re.sub(r"\((\d+)\)", repl, raw_text)


def strip_markers(text):
    """Elimina marcadores tipo [41] y [/41] del texto."""
    return re.sub(r"\[/?(\d+)\]", "", text)


def _strip_option_letter(opt):
    """Elimina el prefijo de letra 'A) ', 'B) ', etc. de una opción."""
    return re.sub(r"^[A-Da-d]\)\s*", "", str(opt))


def format_options_dict(opts):
    """
    Convierte un dict de opciones {A: 'texto', B: 'texto', ...}
    en una lista de cadenas "A) texto".
    """
    return [f"{k}) {v}" for k, v in sorted(opts.items())]


def format_options_list(opts):
    """
    Devuelve la lista de opciones tal cual (ya vienen como "A) texto").
    Si es dict, convierte primero.
    """
    if isinstance(opts, dict):
        return format_options_dict(opts)
    return list(opts)


# ─── funciones por diapositiva ──────────────────────────────────────────────

def set_nombre_test(slide, test_num):
    set_text(slide, "NombreTest", f"Test {test_num}")


def fill_slide_reading1(slide, data, test_num):
    """Diapositiva 2: Reading 1 — Multiple-choice cloze."""
    set_nombre_test(slide, test_num)
    content = data.get("content", {})
    example = content.get("example", {})
    raw_text = content.get("text", "")
    questions = content.get("questions", [])

    set_text(slide, "titulo", data.get("title", ""))
    set_text(slide, "texto", format_parrafos(text_con_huecos(raw_text, questions, example)))

    # Forma "example": "0     opción     opción     opción     opción" (sin letra)
    opts = format_options_list(example.get("options", []))
    ex_num = example.get("number", 0)
    clean_opts = [_strip_option_letter(o) for o in opts]
    parts = [str(ex_num)] + clean_opts
    set_text(slide, "example", "     ".join(parts))

    # Formas A, B, C, D: rellenar de negro la respuesta correcta del example
    correct = example.get("correct", "").upper()
    for letter in ("A", "B", "C", "D"):
        if letter == correct:
            fill_black_shape(slide, letter)
        else:
            clear_shape_fill(slide, letter)

    # Tabla "options": una fila por pregunta (solo texto de la opción, sin letra)
    options_table = get_table(slide, "options")
    for row_idx, q in enumerate(questions):
        q_opts = format_options_list(q.get("options", []))
        fill_table_cell(options_table, row_idx, 0, str(q.get("number", "")))
        col_map = {0: 2, 1: 4, 2: 6, 3: 8}
        for opt_idx, opt in enumerate(q_opts):
            if opt_idx in col_map:
                fill_table_cell(options_table, row_idx, col_map[opt_idx], _strip_option_letter(opt))


def fill_slide_reading2(slide, data, test_num):
    """Diapositiva 3: Reading 2 — Open cloze."""
    set_nombre_test(slide, test_num)
    content = data.get("content", {})
    example = content.get("example", {})
    raw_text = content.get("text", "")
    questions = content.get("questions", [])

    set_text(slide, "titulo", data.get("title", ""))
    set_text(slide, "texto", format_parrafos(text_con_huecos(raw_text, questions, example)))

    # Tabla "example": cada celda contiene una letra de la palabra correcta (mayúscula y negrita)
    correct_word = str(example.get("correct", ""))
    table = get_table(slide, "example")
    for i, letter in enumerate(correct_word):
        fill_table_cell(table, 0, i, letter.upper(), bold=True)


def fill_slide_reading3(slide, data, test_num):
    """Diapositiva 4: Reading 3 — Word formation."""
    set_nombre_test(slide, test_num)
    content = data.get("content", {})
    example = content.get("example", {})
    questions = content.get("questions", [])
    raw_text = content.get("text", "")

    set_text(slide, "titulo", data.get("title", ""))
    set_text(slide, "texto", format_parrafos(text_con_huecos_wf(raw_text, questions, example)))

    # Tabla "example": cada celda contiene una letra de la palabra correcta (negrita)
    correct_word = str(example.get("correct", ""))
    table = get_table(slide, "example")
    for i, letter in enumerate(correct_word):
        fill_table_cell(table, 0, i, letter, bold=True)


def fill_slide_reading4(slide, data, test_num):
    """Diapositiva 5: Reading 4 — Key word transformations."""
    set_nombre_test(slide, test_num)
    content = data.get("content", {})
    questions = content.get("questions", [])

    # Tabla "example": cada celda contiene una letra de la respuesta, empezando
    # en la columna 0 (igual que la tabla "example" de la diapositiva 3).
    example = content.get("example", {})
    if example:
        ex_answer = str(example.get("answer", example.get("correct", ""))).upper()
        table = get_table(slide, "example")
        for i, letter in enumerate(ex_answer):
            fill_table_cell(table, 0, i, letter)

    # Cada pregunta en Frase{i}, key{i} y Refrase{i} (i = posición 1-6)
    for i, q in enumerate(questions, start=1):
        key = q.get("keyWord", q.get("key", ""))
        first = q.get("firstSentence", "")
        before = q.get("beforeGap", "")
        after = q.get("afterGap", "")
        gap_line = f"{before} .................... {after}"
        set_text(slide, f"Frase{i}", first)
        set_text(slide, f"key{i}", key)
        set_text(slide, f"Refrase{i}", gap_line)


def fill_slide_reading5_texto(slide, data, test_num):
    """Diapositiva 6: Reading 5 — título y texto del artículo."""
    set_nombre_test(slide, test_num)
    content = data.get("content", {})
    title = content.get("title", data.get("title", ""))
    raw_text = strip_markers(content.get("text", ""))
    set_text(slide, "titulo", title)
    set_text(slide, "texto", format_parrafos(raw_text))


def fill_slide_reading5_preguntas(slide, data, test_num):
    """Diapositiva 7: Reading 5 — preguntas 31–36 con opciones."""
    set_nombre_test(slide, test_num)
    content = data.get("content", {})
    questions = content.get("questions", [])
    all_runs = []
    for q in questions:
        num = q.get("number", "")
        question_text = q.get("question", "")
        opts = format_options_list(q.get("options", []))
        # Texto enriquecido: número + pregunta, luego opciones con letra en negrita
        all_runs.append({"text": f"{num}. {question_text}\n", "bold": False})
        for opt in opts:
            # opt tiene formato "A) texto..."
            if opt and len(opt) >= 2:
                letter = opt[0]
                rest = opt[1:]
                all_runs.append({"text": f"   {letter}", "bold": True})
                all_runs.append({"text": f"{rest}\n", "bold": False})
    set_shape_text_rich(slide, "preguntas", all_runs)


def fill_slide_reading6_texto(slide, data, test_num):
    """Diapositiva 8: Reading 6 — título y texto con párrafos A–D."""
    set_nombre_test(slide, test_num)
    content = data.get("content", {})
    title = content.get("title", data.get("title", ""))
    texts = content.get("texts", {})
    set_text(slide, "titulo", title)
    runs = []
    for letter in ("A", "B", "C", "D"):
        raw = strip_markers(texts.get(letter, "")).strip()
        runs.append({"text": f"{letter})", "bold": True})
        runs.append({"text": f" {raw}\n\n", "bold": False})
    set_shape_text_rich(slide, "texto", runs)


def fill_slide_reading6_tablas(slide, data, test_num):
    """Diapositiva 9: Reading 6 — tablas de preguntas 37–40."""
    set_nombre_test(slide, test_num)
    content = data.get("content", {})
    texts = content.get("texts", {})
    questions = texts.get("questions", content.get("questions", []))
    if not isinstance(questions, list):
        questions = []
    for q in questions:
        num = q.get("number", "")
        question_text = q.get("question", "")
        table = get_table(slide, str(num))
        fill_table_cell(table, 0, 0, question_text)


def fill_slide_reading7_texto(slide, data, test_num):
    """Diapositiva 10: Reading 7 — texto con numeración entre párrafos (sin título)."""
    set_nombre_test(slide, test_num)
    content = data.get("content", {})
    raw_text = strip_markers(content.get("text", ""))
    set_text(slide, "texto", format_parrafos(raw_text))


def fill_slide_reading7_parrafos(slide, data, test_num):
    """Diapositiva 11: Reading 7 — párrafos A–G en dos formas."""
    set_nombre_test(slide, test_num)
    content = data.get("content", {})
    paragraphs = content.get("paragraphs", {})

    runs1 = []
    for letter in ("A", "B", "C", "D"):
        raw = strip_markers(paragraphs.get(letter, "")).strip()
        runs1.append({"text": f"{letter}\n", "bold": True})
        runs1.append({"text": f"{raw}\n", "bold": False})
    set_shape_text_rich(slide, "parrafos1", runs1)

    runs2 = []
    for letter in ("E", "F", "G"):
        raw = strip_markers(paragraphs.get(letter, "")).strip()
        runs2.append({"text": f"{letter}\n", "bold": True})
        runs2.append({"text": f"{raw}\n", "bold": False})
    set_shape_text_rich(slide, "parrafos2", runs2)


def fill_slide_reading8_tablas(slide, data, test_num):
    """Diapositiva 12: Reading 8 — tablas de preguntas 47–56."""
    set_nombre_test(slide, test_num)
    content = data.get("content", {})
    texts = content.get("texts", {})
    questions = texts.get("questions", content.get("questions", []))
    if not isinstance(questions, list):
        questions = []
    for q in questions:
        num = q.get("number", "")
        question_text = q.get("question", "")
        table = get_table(slide, str(num))
        fill_table_cell(table, 0, 0, question_text)


def fill_slide_reading8_texto(slide, data, test_num):
    """Diapositiva 13: Reading 8 — título y texto con párrafos A–D."""
    set_nombre_test(slide, test_num)
    content = data.get("content", {})
    title = content.get("title", data.get("title", ""))
    texts = content.get("texts", {})
    set_text(slide, "titulo", title)
    runs = []
    for letter in ("A", "B", "C", "D"):
        # Quitar marcadores y encabezados ### al inicio (el JSON puede incluir
        # títulos de estilo Markdown que no deben mostrarse en la diapositiva)
        raw = strip_markers(texts.get(letter, "")).strip()
        raw = re.sub(r"^###\s*", "", raw)
        runs.append({"text": f"{letter})", "bold": True})
        runs.append({"text": f" {raw}\n\n", "bold": False})
    set_shape_text_rich(slide, "texto", runs)


def fill_slide_writing1(slide, data, test_num):
    """Diapositiva 14: Writing 1 — essay."""
    set_nombre_test(slide, test_num)
    content = data.get("content", {})
    notes = content.get("notes", {})

    set_text(slide, "description", data.get("description", ""))
    set_text(slide, "question", content.get("question", ""))

    methods = notes.get("methods", notes.get("aspects", []))
    set_text(slide, "challenges", "\n".join(f"• {m}" for m in methods))

    opinions = notes.get("opinions", [])
    set_text(slide, "opinions", "\n".join(f'"{o}"' for o in opinions))


def fill_slide_writing2(slide, data, test_num):
    """Diapositiva 15: Writing 2 — opciones (forma "options")."""
    set_nombre_test(slide, test_num)
    content = data.get("content", {})
    tasks = content.get("tasks", [])
    all_text = ""
    for task in tasks[:3]:
        title = task.get("title", "")
        prompt = task.get("prompt", task.get("description", ""))
        all_text += f"{title}\n{prompt}\n\n"
    set_text(slide, "options", all_text.rstrip())


def fill_slide_listening1(slide, data, test_num):
    """Diapositiva 16: Listening 1 — preguntas 1–6 y contextos."""
    set_nombre_test(slide, test_num)
    extracts = data.get("extracts", [])
    shape_names = []
    for ctx_idx, extract in enumerate(extracts, start=1):
        context_text = extract.get("context", "")
        combined = context_text + "\n"
        for q in extract.get("questions", []):
            num = q.get("number", "")
            question_text = q.get("question", "")
            opts = format_options_list(q.get("options", []))
            combined += f"{num}. {question_text}\n" + "\n".join(f"   {o}" for o in opts) + "\n\n"
        set_text(slide, f"context{ctx_idx}", combined.rstrip())
        shape_names.append(f"context{ctx_idx}")
    redistribute_shapes_vertically(slide, shape_names)


def fill_slide_listening2(slide, data, test_num):
    """Diapositiva 17: Listening 2 — sentence completion."""
    set_nombre_test(slide, test_num)
    set_text(slide, "description", data.get("instructions", data.get("description", "")))
    questions = []
    for extract in data.get("extracts", []):
        questions.extend(extract.get("questions", []))
    # Sin prefijo de número; añadir GAP_DOTS tras cada marcador (N) en la pregunta.
    # El patrón (\d+)(\s*\.)? captura también el punto final opcional tras el marcador.
    lines = [
        re.sub(r"\((\d+)\)(\s*\.)?", rf"(\1) {GAP_DOTS}", q.get("question", "").strip())
        for q in questions
    ]
    set_text(slide, "texto", "\n".join(lines))


def fill_slide_listening3(slide, data, test_num):
    """Diapositiva 18: Listening 3 — preguntas 15–20 con opciones."""
    set_nombre_test(slide, test_num)
    set_text(slide, "description", data.get("instructions", data.get("description", "")))
    all_text = ""
    for extract in data.get("extracts", []):
        for q in extract.get("questions", []):
            num = q.get("number", "")
            question_text = q.get("question", "")
            opts = format_options_list(q.get("options", []))
            all_text += f"{num}. {question_text}\n" + "\n".join(f"   {o}" for o in opts) + "\n\n"
    set_text(slide, "preguntas", all_text.rstrip())


def fill_slide_listening4(slide, data, test_num):
    """Diapositiva 19: Listening 4 — multiple matching."""
    set_nombre_test(slide, test_num)
    set_text(slide, "description", data.get("description", ""))
    content = data.get("content", {})
    task1 = content.get("task1", {})
    task2 = content.get("task2", {})

    set_text(slide, "instruction1", task1.get("instruction", ""))
    set_text(slide, "instruction2", task2.get("instruction", ""))

    opts1 = task1.get("options", {})
    opts2 = task2.get("options", {})
    set_text(slide, "options1", "\n".join(format_options_dict(opts1) if isinstance(opts1, dict) else opts1))
    set_text(slide, "options2", "\n".join(format_options_dict(opts2) if isinstance(opts2, dict) else opts2))


def _build_answer_pairs(data, part_idx):
    """
    Devuelve una lista de tuplas (num, respuesta) para una parte del examen.
    Incluye el ejemplo (0) si existe.
    """
    content = data.get("content", {})
    pairs = []

    example = content.get("example", {})
    if example:
        correct = example.get("correct", "")
        pairs.append((str(example.get("number", 0)), str(correct)))

    if part_idx == 4:
        # Reading 4: key word transformations — respuesta desde routes
        for q in content.get("questions", []):
            num = q.get("number", "")
            routes = q.get("routes", [])
            if routes:
                route_strs = [
                    f"{r.get('p1', '')} {r.get('p2', '')}".strip() for r in routes
                ]
                correct = " / ".join(route_strs)
            else:
                correct = q.get("correct", q.get("answer", ""))
            pairs.append((str(num), str(correct)))
    else:
        # Algunos ejercicios (p.ej. reading6/reading8 de C1) guardan las
        # preguntas bajo content.texts.questions en lugar de content.questions
        questions = content.get("questions", [])
        if not questions:
            texts = content.get("texts", {})
            if isinstance(texts, dict):
                questions = texts.get("questions", [])
        for q in questions:
            num = q.get("number", "")
            correct = q.get("correct", q.get("answer", ""))
            if isinstance(correct, list):
                correct = " / ".join(correct)
            pairs.append((str(num), str(correct)))

    return pairs


def _fill_answer_table(slide, shape_name, pairs):
    """
    Rellena una tabla de respuestas en horizontal.
    Columnas impares (índice par 0,2,4…) = número de pregunta;
    columnas pares (índice impar 1,3,5…) = respuesta.
    """
    table = get_table(slide, shape_name)
    for col_offset, (num, answer) in enumerate(pairs):
        fill_table_cell(table, 0, col_offset * 2, num)
        fill_table_cell(table, 0, col_offset * 2 + 1, answer)


def _format_answers_horizontal(pairs, cols=4):
    """
    Formatea pares (num, respuesta) en un grid horizontal con columnas de ancho fijo.
    """
    if not pairs:
        return ""
    col_width = max((len(f"{n}. {a}") for n, a in pairs), default=10) + 4
    rows = []
    for i in range(0, len(pairs), cols):
        row_pairs = pairs[i:i + cols]
        rows.append("".join(f"{n}. {a}".ljust(col_width) for n, a in row_pairs).rstrip())
    return "\n".join(rows)


def _format_answers_vertical(pairs):
    """
    Formatea pares (num, respuesta) con cada par en su propia línea.
    Evita que respuestas largas se corten en mitad de línea.
    """
    return "\n".join(f"{n}. {a}" for n, a in pairs)


def fill_slide_answer_reading(slide, test_num, reading_data):
    """Diapositiva 20: Answer Key — Reading & Use of English (formas 1–8)."""
    set_nombre_test(slide, test_num)
    # Partes 1, 5, 6, 7, 8 son tablas; partes 2, 3, 4 son cuadros de texto
    TABLE_PARTS = {1, 5, 6, 7, 8}
    for part_idx, data in enumerate(reading_data, start=1):
        if data is None:
            continue
        pairs = _build_answer_pairs(data, part_idx)
        # Excluir el ejemplo (0) de la hoja de respuestas
        pairs = [(n, a) for n, a in pairs if n != "0"]
        if part_idx in TABLE_PARTS:
            _fill_answer_table(slide, str(part_idx), pairs)
        else:
            # Parts 2, 3, 4: una respuesta por línea para evitar cortes
            set_text(slide, str(part_idx), _format_answers_vertical(pairs))


def fill_slide_answer_listening(slide, test_num, listening_data):
    """Diapositiva 21: Answer Key — Listening (formas 1–4)."""
    set_nombre_test(slide, test_num)
    # Partes 1, 3, 4 son tablas; parte 2 es cuadro de texto
    TABLE_PARTS = {1, 3, 4}
    for part_idx, data in enumerate(listening_data, start=1):
        if data is None:
            continue
        pairs = []
        if part_idx == 4:
            content = data.get("content", {})
            for task_key in ("task1", "task2"):
                task = content.get(task_key, {})
                for q in task.get("questions", []):
                    num = q.get("number", "")
                    correct = q.get("correct", q.get("answer", ""))
                    pairs.append((str(num), str(correct)))
        else:
            for extract in data.get("extracts", []):
                for q in extract.get("questions", []):
                    num = q.get("number", "")
                    correct = q.get("answer", q.get("correct", ""))
                    if isinstance(correct, list):
                        correct = " / ".join(correct)
                    pairs.append((str(num), str(correct)))
        # Excluir el ejemplo (0)
        pairs = [(n, a) for n, a in pairs if n != "0"]
        if part_idx in TABLE_PARTS:
            _fill_answer_table(slide, str(part_idx), pairs)
        else:
            # Part 2: una respuesta por línea para evitar cortes
            set_text(slide, str(part_idx), _format_answers_vertical(pairs))


# ─── carga de datos ─────────────────────────────────────────────────────────

def load_json(path):
    """Carga un JSON desde disco. Devuelve None si el archivo no existe."""
    if not os.path.exists(path):
        return None
    with open(path, encoding="utf-8") as f:
        return json.load(f)


# ─── main ────────────────────────────────────────────────────────────────────

def main():
    if len(sys.argv) < 2:
        print("Uso: python generar_examen.py <ruta_test> [plantilla.pptx]")
        sys.exit(1)

    test_dir = sys.argv[1]
    template_path = sys.argv[2] if len(sys.argv) > 2 else TEMPLATE_DEFAULT

    test_name = os.path.basename(test_dir.rstrip("/\\"))
    test_num = re.sub(r"[^0-9]", "", test_name) or "?"

    print(f"Procesando Test {test_num} desde: {test_dir}")

    def load(name):
        return load_json(os.path.join(test_dir, name))

    reading = [load(f"reading{i}.json") for i in range(1, 9)]
    listening = [load(f"listening{i}.json") for i in range(1, 5)]
    writing = [load("writing1.json"), load("writing2.json")]

    prs = Presentation(template_path)
    slides = prs.slides

    def slide(idx):
        return slides[idx] if idx < len(slides) else None

    print("  Procesando diapositiva 1 (portada)...")
    sl = slide(0)
    if sl:
        set_nombre_test(sl, test_num)

    print("  Procesando diapositiva 2 (Reading 1)...")
    sl = slide(1)
    if sl and reading[0]:
        fill_slide_reading1(sl, reading[0], test_num)

    print("  Procesando diapositiva 3 (Reading 2)...")
    sl = slide(2)
    if sl and reading[1]:
        fill_slide_reading2(sl, reading[1], test_num)

    print("  Procesando diapositiva 4 (Reading 3)...")
    sl = slide(3)
    if sl and reading[2]:
        fill_slide_reading3(sl, reading[2], test_num)

    print("  Procesando diapositiva 5 (Reading 4)...")
    sl = slide(4)
    if sl and reading[3]:
        fill_slide_reading4(sl, reading[3], test_num)

    print("  Procesando diapositiva 6 (Reading 5 - texto)...")
    sl = slide(5)
    if sl and reading[4]:
        fill_slide_reading5_texto(sl, reading[4], test_num)

    print("  Procesando diapositiva 7 (Reading 5 - preguntas)...")
    sl = slide(6)
    if sl and reading[4]:
        fill_slide_reading5_preguntas(sl, reading[4], test_num)

    print("  Procesando diapositiva 8 (Reading 6 - texto)...")
    sl = slide(7)
    if sl and reading[5]:
        fill_slide_reading6_texto(sl, reading[5], test_num)

    print("  Procesando diapositiva 9 (Reading 6 - tablas)...")
    sl = slide(8)
    if sl and reading[5]:
        fill_slide_reading6_tablas(sl, reading[5], test_num)

    print("  Procesando diapositiva 10 (Reading 7 - texto)...")
    sl = slide(9)
    if sl and reading[6]:
        fill_slide_reading7_texto(sl, reading[6], test_num)

    print("  Procesando diapositiva 11 (Reading 7 - párrafos)...")
    sl = slide(10)
    if sl and reading[6]:
        fill_slide_reading7_parrafos(sl, reading[6], test_num)

    print("  Procesando diapositiva 12 (Reading 8 - tablas)...")
    sl = slide(11)
    if sl and reading[7]:
        fill_slide_reading8_tablas(sl, reading[7], test_num)

    print("  Procesando diapositiva 13 (Reading 8 - texto)...")
    sl = slide(12)
    if sl and reading[7]:
        fill_slide_reading8_texto(sl, reading[7], test_num)

    print("  Procesando diapositiva 14 (Writing 1)...")
    sl = slide(13)
    if sl and writing[0]:
        fill_slide_writing1(sl, writing[0], test_num)

    print("  Procesando diapositiva 15 (Writing 2)...")
    sl = slide(14)
    if sl and writing[1]:
        fill_slide_writing2(sl, writing[1], test_num)

    print("  Procesando diapositiva 16 (Listening 1)...")
    sl = slide(15)
    if sl and listening[0]:
        fill_slide_listening1(sl, listening[0], test_num)

    print("  Procesando diapositiva 17 (Listening 2)...")
    sl = slide(16)
    if sl and listening[1]:
        fill_slide_listening2(sl, listening[1], test_num)

    print("  Procesando diapositiva 18 (Listening 3)...")
    sl = slide(17)
    if sl and listening[2]:
        fill_slide_listening3(sl, listening[2], test_num)

    print("  Procesando diapositiva 19 (Listening 4)...")
    sl = slide(18)
    if sl and listening[3]:
        fill_slide_listening4(sl, listening[3], test_num)

    print("  Procesando diapositiva 20 (Answer Key Reading)...")
    sl = slide(19)
    if sl:
        fill_slide_answer_reading(sl, test_num, reading)

    print("  Procesando diapositiva 21 (Answer Key Listening)...")
    sl = slide(20)
    if sl:
        fill_slide_answer_listening(sl, test_num, listening)

    output_path = os.path.join(test_dir, f"Test{test_num}_examen.pptx")
    prs.save(output_path)
    print(f"✅ Guardado en: {output_path}")


if __name__ == "__main__":
    main()
